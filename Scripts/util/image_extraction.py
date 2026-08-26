import hashlib

from pptx.enum.shapes import MSO_SHAPE_TYPE

from Scripts.util.vision_utils import describe_image

# Filters out logos/icons/bullets. PPTX shape size is in EMU (1 inch).
# PDF images are filtered by on-page *display* size in points (72pt = 1 inch),
# not the embedded image's intrinsic pixel resolution -- a small logo can wrap
# a high-resolution source image, and pixel dimensions alone would miss it.
MIN_PPTX_IMAGE_EMU = 914400
MIN_PDF_IMAGE_DIM_PT = 150


def _qualifying_pdf_image_xrefs(page):
    """xrefs of images on this page whose largest on-page placement clears
    MIN_PDF_IMAGE_DIM_PT in both dimensions."""
    xrefs = []
    for img in page.get_images(full=True):
        xref = img[0]
        rects = page.get_image_rects(xref)
        if any(r.width >= MIN_PDF_IMAGE_DIM_PT and r.height >= MIN_PDF_IMAGE_DIM_PT for r in rects):
            xrefs.append(xref)
    return xrefs


def describe_pdf_page_if_image_only(page, seen_xrefs=None):
    """Given a fitz Page with no extractable text, return a vision description
    of the rendered page if it contains a sufficiently large embedded image,
    else None (a truly blank/divider page, or one with only a small logo)."""
    qualifying = _qualifying_pdf_image_xrefs(page)
    if not qualifying:
        return None
    try:
        pix = page.get_pixmap(dpi=150)
        description = describe_image(pix.tobytes("png"))
        if seen_xrefs is not None:
            seen_xrefs.update(qualifying)
        return description
    except Exception as e:
        print(f"Vision description failed for page {page.number + 1}: {e}")
        return None


def describe_pdf_page_images(page, seen_xrefs=None):
    """Return vision descriptions for sufficiently large embedded images on a
    page that already has its own extractable text, so figures alongside
    body text/bullets aren't ignored. seen_xrefs, if given, is a set shared
    across the whole document so a repeated logo/banner image is only
    described once."""
    descriptions = []
    for xref in _qualifying_pdf_image_xrefs(page):
        if seen_xrefs is not None:
            if xref in seen_xrefs:
                continue
            seen_xrefs.add(xref)
        try:
            base_image = page.parent.extract_image(xref)
            mime_type = f"image/{base_image['ext']}"
            descriptions.append(describe_image(base_image["image"], mime_type))
        except Exception as e:
            print(f"Vision description failed for an embedded image: {e}")
    return descriptions


def describe_pptx_slide_images(slide, seen_image_hashes=None):
    """Return a list of vision descriptions for sufficiently large pictures
    on a slide (blank or text-bearing). seen_image_hashes, if given, is a set
    shared across the whole presentation so a repeated logo/banner image is
    only described once."""
    descriptions = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if shape.width < MIN_PPTX_IMAGE_EMU or shape.height < MIN_PPTX_IMAGE_EMU:
            continue
        try:
            image = shape.image
            if seen_image_hashes is not None:
                image_hash = hashlib.sha256(image.blob).hexdigest()
                if image_hash in seen_image_hashes:
                    continue
                seen_image_hashes.add(image_hash)
            descriptions.append(describe_image(image.blob, image.content_type))
        except Exception as e:
            print(f"Vision description failed for a slide image: {e}")
    return descriptions


def get_pptx_slide_notes(slide):
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text
        if notes and notes.strip():
            return notes.strip()
    return None
